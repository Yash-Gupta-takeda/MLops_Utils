from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Takeda colors
TAKEDA_RED = RGBColor(230, 0, 18)
TAKEDA_BLUE = RGBColor(26, 26, 102)
TAKEDA_GREY = RGBColor(246, 246, 246)
TAKEDA_WHITE = RGBColor(255, 255, 255)

# Helper for adding a colored title
def add_title(slide, text, color=TAKEDA_BLUE):
    title = slide.shapes.title
    title.text = text
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = color

def set_background(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_bullet_slide(prs, title, bullets, bullet_color=TAKEDA_BLUE):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    set_background(slide, TAKEDA_WHITE)
    add_title(slide, title, bullet_color)
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()
    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(20)
    return slide

def add_table_slide(prs, title, headers, data, bg_color=TAKEDA_BLUE):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    set_background(slide, TAKEDA_WHITE)
    add_title(slide, title, bg_color)
    rows = len(data) + 1
    cols = len(headers)
    left = Inches(0.5)
    top = Inches(2)
    width = Inches(9)
    height = Inches(0.8)+(0.4*rows)
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    for col, header in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg_color
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(16)
            p.font.color.rgb = TAKEDA_WHITE

    for row_idx, row_data in enumerate(data):
        for col_idx, value in enumerate(row_data):
            cell = table.cell(row_idx+1, col_idx)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = TAKEDA_GREY
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(16)

    return slide


slides_content = [
    {
    "type": "title",
    "title": "MLOps Utility for Real-Time Data Drift & Quality Monitoring",
    "body": [
        "Continuous, platform-agnostic monitoring to detect data shifts, assess quality, and track model performance across diverse environments for reliable ML operations.",
        "Yash"
    ]
},
{
    "type": "bullets",
    "title": "Continuous Monitoring for Reliable Machine Learning",
    "body": [
        "Automates detection of data drift, ensuring model relevance in production",
        "Monitors data quality metrics continuously to maintain integrity",
        "Tracks comprehensive model performance metrics to assess accuracy and fairness",
        "Supports multiple data formats and integrates with MLflow, pickle artifacts, and data pipelines",
        "Enables real-time observability and governance for proactive feature shift detection",
        "Exports consolidated monitoring reports facilitating transparency and regulatory compliance",
        "Ensures sustained trust and compliance in automated decision-making systems"
    ]
},
{
    "type": "bullets",
    "title": "Critical Challenges Driving ML Model Degradation",
    "body": [
        "Shifts in input data distributions (feature drift) reduce model accuracy",
        "Inconsistencies and errors from complex data pipelines introduce noise",
        "Changes in the relationship between inputs and outputs (concept drift) cause outdated predictions",
        "Manual, fragmented monitoring delays problem detection and response",
        "Model degradation leads to biased outcomes, reduced trust, and increased compliance risks",
        "Organizations require automated, unified, and platform-agnostic monitoring to maintain model integrity"
    ]
},
{
    "type": "bullets",
    "title": "DataDriftSDK: Advanced Data Drift Detection",
    "body": [
        "Calculates detailed data drift metrics using Evidently AI to track input feature shifts over time",
        "Supports multiple data formats: CSV, Parquet, and Spark/Delta tables for versatile ingestion",
        "Generates both per-feature and aggregate dataset drift scores for granular and overall analysis",
        "Configurable statistical methods and drift thresholds via customizable configuration files",
        "Exports results in a platform-agnostic way to local storage or Delta tables, enabling seamless integration",
        "Enables proactive detection of problematic data shifts before model performance degrades"
    ]
},
{
    "type": "bullets",
    "title": "DataQualitySDK Delivers Comprehensive Data Profiling",
    "body": [
        "Calculates essential metrics: missing value ratios, variance, mean, standard deviation, cardinality, and outlier detection",
        "Identifies feature types, distinguishing numerical from categorical data",
        "Captures unique and common values per feature for detailed profiling",
        "Stores results in a consistent schema to standardize data quality reporting",
        "Exports profiles to Delta or CSV formats, enabling seamless integration",
        "Supports data engineers by detecting schema drifts and corruption early, ensuring reliable model inputs"
    ]
},
{
    "type": "bullets",
    "title": "Optimizing Model Performance with ModelMonitoringSDK",
    "body": [
        "Evaluates ML model performance for classification and regression tasks with metrics like Accuracy, F1-Score, RMSE, and R²",
        "Seamlessly integrates with MLflow and Pickle model artifacts to fit diverse ML workflows",
        "Employs SHAP values to deliver clear model explainability and interpret complex predictions",
        "Correlates performance metrics with feature drift detected by DataDriftSDK to uncover impact drivers",
        "Exports results to CSV or Databricks Delta tables for unified tracking of model behavior and data shifts",
        "Enables enhanced interpretability and governance by tracking model metrics alongside detected data drift"
    ]
},
{
    "type": "bullets",
    "title": "Stakeholders Empowered by Monitoring",
    "body": [
        "Data Scientists monitor feature drift and model accuracy in real time to maintain model performance.",
        "MLOps Engineers automate retraining triggers and build governance dashboards for streamlined model management.",
        "Data Engineers ensure pipeline data quality and reliability, validating inputs for robust ML workflows.",
        "Business Analysts translate monitoring results into actionable insights linked to KPIs and business goals.",
        "Compliance Teams verify auditability, fairness, and regulatory documentation to maintain ethical standards.",
        "This cross-functional targeting fosters collaboration and shared responsibility for sustained ML system health."
    ]
},
{
    "type": "bullets",
    "title": "Comprehensive Monitoring Utility Advantages",
    "body": [
        "End-to-end monitoring covering data drift, data quality, and model performance within a single platform",
        "Platform-agnostic support enabling deployment on local, cloud, and Databricks environments",
        "Configurable thresholds that enable governance and alerting without requiring code changes",
        "SHAP-based explainability to enhance model transparency and interpretability",
        "Automated export features centralizing reporting and easing compliance workflows",
        "Lightweight SDK design facilitating seamless integration into existing ML pipelines",
        "Audit trail maintenance critical for compliance and traceability",
        "Extensibility allowing addition of new drift tests or quality metrics to meet evolving needs"
    ]
},
{
    "type": "bullets",
    "title": "Seamless Integration Across ML Workflows",
    "body": [
        "Integrates with popular CI/CD pipelines including Jenkins, GitHub Actions, and Azure DevOps",
        "Supports orchestration tools like Airflow, Databricks Jobs, and Prefect for scheduling monitoring tasks",
        "Version control integration via MLflow ensures consistent experiment tracking and artifact management",
        "Offers visualization through Evidently dashboards and Power BI plugins for insightful, interactive reporting",
        "Ensures consistent, automated monitoring spanning from model training to production deployment"
    ]
},
{
    "type": "bullets",
    "title": "Driving Scalable, Transparent, and Compliant ML Operations",
    "body": [
        "Provides a robust and scalable foundation for high-quality ML operations",
        "Integrates data drift detection, data quality profiling, and model performance tracking",
        "Enables teams to proactively manage model reliability and regulatory compliance",
        "Offers a platform-agnostic, modular SDK design compatible with diverse environments",
        "Supports delivery of trustworthy, efficient ML systems aligning with evolving business needs"
    ]
},
{
    "type": "table",
    "title": "Modular SDKs Simplifying MLOps Monitoring",
    "headers": ["SDK Component", "Purpose", "Outputs"],
    "rows": [
        ["DataDriftSDK", "Detects distribution shifts between training and scoring datasets", "Drift metrics indicating data distribution changes"],
        ["DataQualitySDK", "Profiles key data quality metrics including missing values and outliers", "Data quality reports highlighting anomalies"],
        ["ModelMonitoringSDK", "Evaluates classification and regression models with metrics and SHAP explainability integrating drift insights", "Model performance metrics & explainability integrated with drift data"]
    ]
},
{
    "type": "table",
    "title": "Business Use Cases: Practical Applications Across Industries",
    "headers": ["Business Scenario", "Description", "Benefits"],
    "rows": [
        ["Predictive Sales Models", "Monitors customer and transaction patterns to detect and prevent forecast degradation.", "Preserves forecast accuracy for better sales planning and revenue management."],
        ["Risk and Fraud Detection", "Tracks transactional data drifts to maintain model precision and reliability.", "Enhances fraud prevention, reducing financial losses and reputational damage."],
        ["Customer Churn Models", "Continuously monitors critical features to sustain predictive accuracy over time.", "Supports timely customer retention actions to reduce churn rates."],
        ["Supply Chain Optimization", "Ensures consistency and reliability of logistics data to prevent disruptions.", "Improves operational efficiency and minimizes downtime in supply chains."],
        ["Marketing Personalization", "Monitors campaign and user data quality to maintain targeting effectiveness.", "Increases campaign ROI and customer engagement through reliable personalization."]
    ]
},
{
    "type": "table",
    "title": "Transforming Operations with Automation",
    "headers": ["Task", "Before Automation", "After Automation"],
    "rows": [
        ["Data Drift Detection", "Manual SQL queries and visual analysis", "Automated detection using Evidently AI"],
        ["Quality Validation", "Ad-hoc scripts", "Continuous profiling integrated in monitoring"],
        ["Model Monitoring", "Manual metric tracking with limited environment support", "Unified tracking via MLflow, SHAP, and drift insights"],
        ["Environment Compatibility", "Limited", "Platform-agnostic: supports local, cloud, and Databricks"],
        ["Metric Export & Aggregation", "Manual aggregation", "Automated metric export streamlining workflows"],
        ["Model Reliability & Governance", "Challenging due to manual processes", "Enhanced through automation and unified tracking"]
    ]
},
{
    "type": "table",
    "title": "Platform-Agnostic Design: Flexible Multi-Environment Support",
    "headers": ["Platform", "Support Type", "Data Export Formats", "Integration Mode"],
    "rows": [
        ["Local Systems", "Automated environment detection and I/O handling", "CSV, Parquet", "Smooth integration into existing CI/CD pipelines"],
        ["Databricks with Unity Catalog", "Auto-detects environment for optimized data operations", "Delta tables", "Platform-neutral deployment enabling flexible workflows"],
        ["AWS Redshift", "Environment-aware I/O management", "SQL, CSV", "Easy integration without platform lock-in"],
        ["On-premises & Custom Cloud", "Custom environment detection and adaptable handling", "CSV, Parquet, SQL", "Supports diverse organizational infrastructures"]
    ]
},
{
    "type": "table",
    "title": "Comprehensive Technology Stack & Language Coverage",
    "headers": ["Technology", "Purpose / Usage"],
    "rows": [
        ["Python 3.9+", "Core implementation language"],
        ["Evidently AI", "Drift detection in datasets"],
        ["Pandas, NumPy", "Data processing and manipulation"],
        ["PySpark", "Scalable data workflows"],
        ["SQLAlchemy", "Database interaction and ORM"],
        ["MLflow", "Experiment and model tracking"],
        ["ExplainerDashboard", "SHAP-based model explainability"],
        ["Data Storage Formats", "CSV, Parquet, Delta tables, Redshift"],
        ["Logging", "Python’s logging module for event recording"],
        ["Configuration", "INI files enabling flexible parameter tuning"]
    ]
},
{
    "type": "table",
    "title": "Comprehensive Metrics Tracking for Robust Governance",
    "headers": ["Metric Type", "Source/SDK", "Tracked Elements", "Governance Purpose"],
    "rows": [
        ["Data Drift Metrics", "Evidently AI (DataDriftSDK)", "Shifts in input distributions", "Detects data input variations impacting model validity"],
        ["Data Quality Metrics", "DataQualitySDK", "Missingness, Variance", "Monitors data integrity and consistency"],
        ["Model Performance Metrics", "ModelMonitoringSDK", "Accuracy, Recall, RMSE, SHAP importances", "Evaluates model effectiveness and explains feature impact"],
        ["Operational Metrics", "Local logs / Spark", "Runtime, Frequency", "Tracks execution details supporting operational oversight"],
        ["Record Attributes", "All Systems", "Timestamp, Model Name, Metric Scores, Drifted Feature Counts, Environment Info", "Supports governance dashboards and automated retraining triggers"]
    ]
},
{
    "type": "table",
    "title": "Future Scope: Expanding Monitoring Capabilities",
    "headers": ["Feature Area", "Planned Enhancements", "Value / Benefits"],
    "rows": [
        ["Concept Drift Detection", "Advanced methods: Kolmogorov–Smirnov, PSI, MMD", "Improves detection accuracy for model aging and distribution shifts"],
        ["Real-Time Alerts", "Integration with Slack and Teams for anomaly notifications", "Enhances responsiveness and proactive issue resolution"],
        ["Visualization Dashboard", "Power BI and Evidently UI integration layer", "Delivers intuitive data and drift insights for stakeholders"],
        ["Automated Retraining", "Pipelines triggered by drift thresholds", "Streamlines model maintenance and reduces manual intervention"],
        ["Schema Evolution Monitoring", "Track changes in data schema over time", "Ensures data integrity and compatibility with model inputs"],
        ["Support for Complex Models", "Extend monitoring to time-series and deep learning models", "Broadens applicability to diverse ML operational needs"]
    ]
}
]


from pptx import Presentation
from pptx.util import Pt
# ...plus your helper functions and color definitions from before...

prs = Presentation()

for idx, slide_info in enumerate(slides_content):
    if slide_info["type"] == "title":
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        set_background(slide, TAKEDA_WHITE)
        title_shape = slide.shapes.title
        title_shape.text = slide_info["title"]
        title_shape.text_frame.paragraphs[0].font.size = Pt(36)
        title_shape.text_frame.paragraphs[0].font.bold = True          
        subtitle_shape = slide.placeholders[1]
        subtitle_shape.text = "\n\n".join(slide_info["body"])
        subtitle_shape.text_frame.paragraphs[0].font.size = Pt(20)
        subtitle_shape.text_frame.paragraphs[0].font.color.rgb = TAKEDA_BLUE
    elif slide_info["type"] == "bullets":
        add_bullet_slide(prs, slide_info["title"], slide_info["body"])
    elif slide_info["type"] == "table":
        add_table_slide(prs, slide_info["title"], slide_info["headers"], slide_info["rows"])
    else:
        raise ValueError("Unknown slide type: {}".format(slide_info["type"]))

prs.save("Takeda_MLOps_Monitoring.pptx")
print("PPTX file generated: Takeda_MLOps_Monitoring.pptx")




